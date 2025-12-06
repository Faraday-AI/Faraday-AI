"""
Widget Export Service

Provides functionality to export widget data to various formats:
- PDF
- Word (DOCX)
- Excel (XLSX)
- PowerPoint (PPTX)
- Email (via communication service)
"""

import io
import logging
import re
from typing import Dict, Any, Optional, List
from datetime import datetime
import json

# PDF
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak, Image as ReportLabImage
from reportlab.lib.enums import TA_CENTER, TA_LEFT

# Word
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH

# Excel
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils import get_column_letter

# PowerPoint
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Image handling
from PIL import Image as PILImage
import base64
import requests

logger = logging.getLogger(__name__)


class WidgetExportService:
    """Service for exporting widget data to various formats."""
    
    def __init__(self):
        self.logger = logger
    
    def export_to_pdf(self, widget_data: Dict[str, Any], widget_type: str, widget_title: str) -> bytes:
        """Export widget data to PDF format."""
        try:
            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=letter)
            story = []
            styles = getSampleStyleSheet()
            
            # Title style
            title_style = ParagraphStyle(
                'WidgetTitle',
                parent=styles['Heading1'],
                fontSize=20,
                textColor=colors.HexColor('#1a1a1a'),
                spaceAfter=20,
                alignment=TA_CENTER
            )
            
            # Add title
            story.append(Paragraph(widget_title, title_style))
            story.append(Spacer(1, 0.2*inch))
            
            # Metadata
            export_date = datetime.utcnow().strftime("%Y-%m-%d %H:%M:%S UTC")
            story.append(Paragraph(f"<b>Widget Type:</b> {widget_type}", styles['Normal']))
            story.append(Paragraph(f"<b>Export Date:</b> {export_date}", styles['Normal']))
            story.append(Spacer(1, 0.2*inch))
            
            # Format widget data
            formatted_content = self._format_data_for_pdf(widget_data, widget_type)
            story.extend(formatted_content)
            
            # Build PDF
            doc.build(story)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            self.logger.error(f"Error exporting widget to PDF: {str(e)}")
            raise
    
    def export_to_word(self, widget_data: Dict[str, Any], widget_type: str, widget_title: str) -> bytes:
        """Export widget data to Word (DOCX) format."""
        try:
            doc = Document()
            
            # Title
            title = doc.add_heading(widget_title, 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Metadata
            doc.add_paragraph(f"Widget Type: {widget_type}")
            doc.add_paragraph(f"Export Date: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S UTC')}")
            doc.add_paragraph()  # Blank line
            
            # Format widget data
            self._format_data_for_word(doc, widget_data, widget_type)
            
            # Save to bytes
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            self.logger.error(f"Error exporting widget to Word: {str(e)}")
            raise
    
    def export_to_excel(self, widget_data: Dict[str, Any], widget_type: str, widget_title: str) -> bytes:
        """Export widget data to Excel (XLSX) format."""
        try:
            wb = Workbook()
            ws = wb.active
            ws.title = widget_title[:31]  # Excel sheet name limit
            
            # Header style
            header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF", size=12)
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # Title row
            ws.merge_cells('A1:B1')
            title_cell = ws['A1']
            title_cell.value = widget_title
            title_cell.font = Font(bold=True, size=14)
            title_cell.alignment = Alignment(horizontal='center', vertical='center')
            
            # Metadata
            ws['A2'] = 'Widget Type'
            ws['B2'] = widget_type
            ws['A3'] = 'Export Date'
            ws['B3'] = datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S UTC')
            
            # Format widget data
            self._format_data_for_excel(ws, widget_data, widget_type, start_row=5)
            
            # Auto-adjust column widths
            for column in ws.columns:
                max_length = 0
                column_letter = get_column_letter(column[0].column)
                for cell in column:
                    try:
                        if len(str(cell.value)) > max_length:
                            max_length = len(str(cell.value))
                    except:
                        pass
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width
            
            # Save to bytes
            buffer = io.BytesIO()
            wb.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            self.logger.error(f"Error exporting widget to Excel: {str(e)}")
            raise
    
    def export_to_powerpoint(self, widget_data: Dict[str, Any], widget_type: str, widget_title: str) -> bytes:
        """Export widget data to PowerPoint (PPTX) format."""
        try:
            prs = Presentation()
            prs.slide_width = PptInches(10)
            prs.slide_height = PptInches(7.5)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = widget_title
            subtitle.text = f"{widget_type}\nExported: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S UTC')}"
            
            # Content slides
            self._format_data_for_powerpoint(prs, widget_data, widget_type)
            
            # Save to bytes
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            self.logger.error(f"Error exporting widget to PowerPoint: {str(e)}")
            raise
    
    def _format_data_for_pdf(self, data: Dict[str, Any], widget_type: str) -> List:
        """Format widget data for PDF export."""
        story = []
        styles = getSampleStyleSheet()
        
        # Special handling for lesson plans
        if widget_type == 'lesson-planning' or 'lesson' in widget_type.lower():
            return self._format_lesson_plan_for_pdf(data, story, styles)
        
        # Special handling for tables (rubrics, worksheets)
        if isinstance(data, dict):
            # Check for rubric data
            if 'rubrics' in data or 'rubric' in str(data).lower():
                return self._format_rubric_for_pdf(data, story, styles)
            
            # Check for worksheet data
            if 'worksheets' in data or 'worksheet' in str(data).lower():
                return self._format_worksheet_for_pdf(data, story, styles)
        
        # Generic formatting
        return self._format_generic_data_for_pdf(data, story, styles)
    
    def _format_lesson_plan_for_pdf(self, data: Dict[str, Any], story: List, styles) -> List:
        """Format lesson plan data for PDF."""
        # Title
        if data.get('title'):
            story.append(Paragraph(f"<b>Title:</b> {data['title']}", styles['Heading2']))
            story.append(Spacer(1, 0.1*inch))
        
        # Description
        if data.get('description'):
            story.append(Paragraph("<b>Description:</b>", styles['Heading3']))
            story.append(Paragraph(data['description'], styles['Normal']))
            story.append(Spacer(1, 0.1*inch))
        
        # Objectives
        if data.get('objectives'):
            story.append(Paragraph("<b>Learning Objectives:</b>", styles['Heading3']))
            objectives = data['objectives'] if isinstance(data['objectives'], list) else [data['objectives']]
            for obj in objectives:
                if obj:
                    story.append(Paragraph(f"• {str(obj)}", styles['Normal']))
            story.append(Spacer(1, 0.1*inch))
        
        # Activities
        if data.get('activities'):
            story.append(Paragraph("<b>Activities:</b>", styles['Heading3']))
            activities = data['activities'] if isinstance(data['activities'], list) else [data['activities']]
            for i, activity in enumerate(activities, 1):
                if isinstance(activity, dict):
                    name = activity.get('name', f'Activity {i}')
                    desc = activity.get('description', '')
                    story.append(Paragraph(f"<b>{i}. {name}</b>", styles['Normal']))
                    if desc:
                        story.append(Paragraph(desc, styles['Normal']))
                else:
                    story.append(Paragraph(f"{i}. {str(activity)}", styles['Normal']))
            story.append(Spacer(1, 0.1*inch))
        
        # Worksheets
        if data.get('worksheets'):
            story.append(PageBreak())
            story.append(Paragraph("<b>Worksheets</b>", styles['Heading2']))
            worksheets_text = data['worksheets']
            if isinstance(worksheets_text, str):
                # Parse and format worksheets
                lines = worksheets_text.split('\n')
                for line in lines:
                    if line.strip():
                        story.append(Paragraph(line.strip(), styles['Normal']))
            story.append(Spacer(1, 0.1*inch))
        
        # Rubrics
        if data.get('rubrics'):
            story.append(PageBreak())
            story.append(Paragraph("<b>Grading Rubric</b>", styles['Heading2']))
            rubrics_text = data['rubrics']
            if isinstance(rubrics_text, str):
                # Try to parse as table
                if '|' in rubrics_text:
                    return self._format_rubric_table_for_pdf(rubrics_text, story, styles)
                else:
                    lines = rubrics_text.split('\n')
                    for line in lines:
                        if line.strip():
                            story.append(Paragraph(line.strip(), styles['Normal']))
        
        return story
    
    def _format_rubric_for_pdf(self, data: Dict[str, Any], story: List, styles) -> List:
        """Format rubric data for PDF."""
        rubrics_text = data.get('rubrics', '')
        if isinstance(rubrics_text, str) and '|' in rubrics_text:
            return self._format_rubric_table_for_pdf(rubrics_text, story, styles)
        return self._format_generic_data_for_pdf(data, story, styles)
    
    def _format_rubric_table_for_pdf(self, rubrics_text: str, story: List, styles) -> List:
        """Format rubric table for PDF."""
        lines = [line.strip() for line in rubrics_text.split('\n') if line.strip() and '|' in line]
        
        if not lines:
            return story
        
        # Parse table rows
        table_data = []
        for line in lines:
            # Skip separator rows
            if re.match(r'^\|[\s\-:]+\|', line):
                continue
            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
            if cells and len(cells) >= 2:
                # Skip header-like rows
                if cells[0].lower() == 'criteria' or 'excellent' in cells[0].lower():
                    if not table_data:  # First row is header
                        table_data.append(cells)
                    continue
                table_data.append(cells)
        
        if table_data:
            # Create PDF table
            table = Table(table_data, colWidths=[1.5*inch] + [1.2*inch] * (len(table_data[0]) - 1))
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, 0), 10),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                ('GRID', (0, 0), (-1, -1), 1, colors.black),
                ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ]))
            story.append(table)
        
        return story
    
    def _format_worksheet_for_pdf(self, data: Dict[str, Any], story: List, styles) -> List:
        """Format worksheet data for PDF."""
        worksheets_text = data.get('worksheets', '')
        if isinstance(worksheets_text, str):
            lines = worksheets_text.split('\n')
            for line in lines:
                if line.strip():
                    # Format questions
                    if re.match(r'^\d+[\.\)]\s+', line):
                        story.append(Paragraph(f"<b>{line.strip()}</b>", styles['Normal']))
                    # Format answer keys
                    elif 'correct answer' in line.lower() or 'answer:' in line.lower():
                        story.append(Paragraph(f"<i>{line.strip()}</i>", styles['Normal']))
                    else:
                        story.append(Paragraph(line.strip(), styles['Normal']))
        return story
    
    def _format_generic_data_for_pdf(self, data: Dict[str, Any], story: List, styles) -> List:
        """Format generic widget data for PDF."""
        if isinstance(data, dict):
            for key, value in data.items():
                if value and key not in ['id', 'widget_id', 'created_at', 'updated_at']:
                    if isinstance(value, (list, dict)):
                        story.append(Paragraph(f"<b>{key.replace('_', ' ').title()}:</b>", styles['Heading3']))
                        story.append(Paragraph(json.dumps(value, indent=2, default=str), styles['Code']))
                    else:
                        story.append(Paragraph(f"<b>{key.replace('_', ' ').title()}:</b> {str(value)}", styles['Normal']))
                    story.append(Spacer(1, 0.05*inch))
        elif isinstance(data, list):
            for i, item in enumerate(data, 1):
                story.append(Paragraph(f"<b>Item {i}:</b>", styles['Heading3']))
                if isinstance(item, dict):
                    for key, value in item.items():
                        story.append(Paragraph(f"{key}: {str(value)}", styles['Normal']))
                else:
                    story.append(Paragraph(str(item), styles['Normal']))
                story.append(Spacer(1, 0.1*inch))
        else:
            story.append(Paragraph(str(data), styles['Normal']))
        
        return story
    
    def _format_data_for_word(self, doc: Document, data: Dict[str, Any], widget_type: str):
        """Format widget data for Word export."""
        # Similar logic to PDF but using docx API
        if widget_type == 'lesson-planning' or 'lesson' in widget_type.lower():
            self._format_lesson_plan_for_word(doc, data)
        elif 'rubric' in str(data).lower():
            self._format_rubric_for_word(doc, data)
        elif 'worksheet' in str(data).lower():
            self._format_worksheet_for_word(doc, data)
        else:
            self._format_generic_data_for_word(doc, data)
    
    def _format_lesson_plan_for_word(self, doc: Document, data: Dict[str, Any]):
        """Format lesson plan for Word."""
        if data.get('title'):
            doc.add_heading(data['title'], level=1)
        if data.get('description'):
            doc.add_heading('Description', level=2)
            doc.add_paragraph(data['description'])
        if data.get('objectives'):
            doc.add_heading('Learning Objectives', level=2)
            objectives = data['objectives'] if isinstance(data['objectives'], list) else [data['objectives']]
            for obj in objectives:
                if obj:
                    doc.add_paragraph(str(obj), style='List Bullet')
        if data.get('activities'):
            doc.add_heading('Activities', level=2)
            activities = data['activities'] if isinstance(data['activities'], list) else [data['activities']]
            for i, activity in enumerate(activities, 1):
                if isinstance(activity, dict):
                    name = activity.get('name', f'Activity {i}')
                    desc = activity.get('description', '')
                    doc.add_paragraph(f"{i}. {name}", style='List Number')
                    if desc:
                        doc.add_paragraph(desc)
                else:
                    doc.add_paragraph(str(activity), style='List Number')
        if data.get('worksheets'):
            doc.add_page_break()
            doc.add_heading('Worksheets', level=1)
            worksheets_text = data['worksheets']
            if isinstance(worksheets_text, str):
                for line in worksheets_text.split('\n'):
                    if line.strip():
                        doc.add_paragraph(line.strip())
        if data.get('rubrics'):
            doc.add_page_break()
            doc.add_heading('Grading Rubric', level=1)
            rubrics_text = data['rubrics']
            if isinstance(rubrics_text, str):
                if '|' in rubrics_text:
                    self._format_rubric_table_for_word(doc, rubrics_text)
                else:
                    for line in rubrics_text.split('\n'):
                        if line.strip():
                            doc.add_paragraph(line.strip())
    
    def _format_rubric_table_for_word(self, doc: Document, rubrics_text: str):
        """Format rubric table for Word."""
        lines = [line.strip() for line in rubrics_text.split('\n') if line.strip() and '|' in line]
        if not lines:
            return
        
        # Parse table
        table_data = []
        for line in lines:
            if re.match(r'^\|[\s\-:]+\|', line):
                continue
            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
            if cells and len(cells) >= 2:
                if cells[0].lower() == 'criteria' or 'excellent' in cells[0].lower():
                    if not table_data:
                        table_data.append(cells)
                    continue
                table_data.append(cells)
        
        if table_data:
            table = doc.add_table(rows=len(table_data), cols=len(table_data[0]))
            table.style = 'Light Grid Accent 1'
            
            for i, row_data in enumerate(table_data):
                for j, cell_data in enumerate(row_data):
                    cell = table.rows[i].cells[j]
                    cell.text = cell_data
                    if i == 0:  # Header row
                        cell.paragraphs[0].runs[0].font.bold = True
    
    def _format_worksheet_for_word(self, doc: Document, data: Dict[str, Any]):
        """Format worksheet for Word."""
        worksheets_text = data.get('worksheets', '')
        if isinstance(worksheets_text, str):
            for line in worksheets_text.split('\n'):
                if line.strip():
                    if re.match(r'^\d+[\.\)]\s+', line):
                        p = doc.add_paragraph(line.strip())
                        p.runs[0].font.bold = True
                    else:
                        doc.add_paragraph(line.strip())
    
    def _format_generic_data_for_word(self, doc: Document, data: Dict[str, Any]):
        """Format generic data for Word."""
        if isinstance(data, dict):
            for key, value in data.items():
                if value and key not in ['id', 'widget_id', 'created_at', 'updated_at']:
                    doc.add_heading(key.replace('_', ' ').title(), level=2)
                    if isinstance(value, (list, dict)):
                        doc.add_paragraph(json.dumps(value, indent=2, default=str))
                    else:
                        doc.add_paragraph(str(value))
        elif isinstance(data, list):
            for i, item in enumerate(data, 1):
                doc.add_heading(f'Item {i}', level=2)
                if isinstance(item, dict):
                    for k, v in item.items():
                        doc.add_paragraph(f"{k}: {v}")
                else:
                    doc.add_paragraph(str(item))
        else:
            doc.add_paragraph(str(data))
    
    def _format_data_for_excel(self, ws, data: Dict[str, Any], widget_type: str, start_row: int = 5):
        """Format widget data for Excel export."""
        current_row = start_row
        
        if widget_type == 'lesson-planning' or 'lesson' in widget_type.lower():
            self._format_lesson_plan_for_excel(ws, data, current_row)
        elif 'rubric' in str(data).lower():
            self._format_rubric_for_excel(ws, data, current_row)
        elif 'worksheet' in str(data).lower():
            self._format_worksheet_for_excel(ws, data, current_row)
        else:
            self._format_generic_data_for_excel(ws, data, current_row)
    
    def _format_lesson_plan_for_excel(self, ws, data: Dict[str, Any], start_row: int):
        """Format lesson plan for Excel."""
        row = start_row
        if data.get('title'):
            ws[f'A{row}'] = 'Title'
            ws[f'B{row}'] = data['title']
            row += 1
        if data.get('description'):
            ws[f'A{row}'] = 'Description'
            ws[f'B{row}'] = data['description']
            row += 1
        if data.get('objectives'):
            ws[f'A{row}'] = 'Objectives'
            objectives = data['objectives'] if isinstance(data['objectives'], list) else [data['objectives']]
            ws[f'B{row}'] = '\n'.join([str(obj) for obj in objectives if obj])
            row += 1
        if data.get('activities'):
            ws[f'A{row}'] = 'Activities'
            activities = data['activities'] if isinstance(data['activities'], list) else [data['activities']]
            activity_text = []
            for i, activity in enumerate(activities, 1):
                if isinstance(activity, dict):
                    name = activity.get('name', f'Activity {i}')
                    desc = activity.get('description', '')
                    activity_text.append(f"{i}. {name}: {desc}")
                else:
                    activity_text.append(f"{i}. {str(activity)}")
            ws[f'B{row}'] = '\n'.join(activity_text)
            row += 1
        if data.get('worksheets'):
            row += 1
            ws[f'A{row}'] = 'Worksheets'
            ws[f'A{row}'].font = Font(bold=True)
            row += 1
            worksheets_text = data['worksheets']
            if isinstance(worksheets_text, str):
                lines = worksheets_text.split('\n')
                for line in lines[:50]:  # Limit to 50 lines
                    if line.strip():
                        ws[f'B{row}'] = line.strip()
                        row += 1
        if data.get('rubrics'):
            row += 1
            ws[f'A{row}'] = 'Rubrics'
            ws[f'A{row}'].font = Font(bold=True)
            row += 1
            rubrics_text = data['rubrics']
            if isinstance(rubrics_text, str) and '|' in rubrics_text:
                self._format_rubric_table_for_excel(ws, rubrics_text, row)
    
    def _format_rubric_table_for_excel(self, ws, rubrics_text: str, start_row: int):
        """Format rubric table for Excel."""
        lines = [line.strip() for line in rubrics_text.split('\n') if line.strip() and '|' in line]
        if not lines:
            return
        
        row = start_row
        table_data = []
        for line in lines:
            if re.match(r'^\|[\s\-:]+\|', line):
                continue
            cells = [cell.strip() for cell in line.split('|') if cell.strip()]
            if cells and len(cells) >= 2:
                if cells[0].lower() == 'criteria' or 'excellent' in cells[0].lower():
                    if not table_data:
                        table_data.append(cells)
                    continue
                table_data.append(cells)
        
        if table_data:
            for i, row_data in enumerate(table_data):
                for j, cell_data in enumerate(row_data):
                    cell = ws.cell(row=row + i, column=j + 1, value=cell_data)
                    if i == 0:  # Header
                        cell.font = Font(bold=True, color="FFFFFF")
                        cell.fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
    
    def _format_worksheet_for_excel(self, ws, data: Dict[str, Any], start_row: int):
        """Format worksheet for Excel."""
        worksheets_text = data.get('worksheets', '')
        if isinstance(worksheets_text, str):
            row = start_row
            lines = worksheets_text.split('\n')
            for line in lines[:100]:  # Limit to 100 lines
                if line.strip():
                    ws[f'A{row}'] = line.strip()
                    row += 1
    
    def _format_generic_data_for_excel(self, ws, data: Dict[str, Any], start_row: int):
        """Format generic data for Excel."""
        row = start_row
        if isinstance(data, dict):
            for key, value in data.items():
                if value and key not in ['id', 'widget_id', 'created_at', 'updated_at']:
                    ws[f'A{row}'] = key.replace('_', ' ').title()
                    if isinstance(value, (list, dict)):
                        ws[f'B{row}'] = json.dumps(value, default=str)
                    else:
                        ws[f'B{row}'] = str(value)
                    row += 1
        elif isinstance(data, list):
            ws[f'A{start_row}'] = 'Items'
            ws[f'A{start_row}'].font = Font(bold=True)
            row = start_row + 1
            for i, item in enumerate(data, 1):
                ws[f'A{row}'] = f"Item {i}"
                if isinstance(item, dict):
                    for j, (k, v) in enumerate(item.items()):
                        ws.cell(row=row, column=j+2, value=f"{k}: {v}")
                else:
                    ws[f'B{row}'] = str(item)
                row += 1
    
    def _format_data_for_powerpoint(self, prs: Presentation, data: Dict[str, Any], widget_type: str):
        """Format widget data for PowerPoint."""
        # Content slide layout
        bullet_slide_layout = prs.slide_layouts[1]
        
        if widget_type == 'lesson-planning' or 'lesson' in widget_type.lower():
            self._format_lesson_plan_for_powerpoint(prs, data, bullet_slide_layout)
        else:
            self._format_generic_data_for_powerpoint(prs, data, bullet_slide_layout)
    
    def _format_lesson_plan_for_powerpoint(self, prs: Presentation, data: Dict[str, Any], layout):
        """Format lesson plan for PowerPoint."""
        if data.get('title'):
            slide = prs.slides.add_slide(layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = data['title']
            if data.get('description'):
                content.text = data['description']
        
        if data.get('objectives'):
            slide = prs.slides.add_slide(layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = "Learning Objectives"
            objectives = data['objectives'] if isinstance(data['objectives'], list) else [data['objectives']]
            tf = content.text_frame
            for obj in objectives:
                if obj:
                    p = tf.add_paragraph()
                    p.text = str(obj)
                    p.level = 0
        
        if data.get('activities'):
            slide = prs.slides.add_slide(layout)
            title = slide.shapes.title
            content = slide.placeholders[1]
            title.text = "Activities"
            activities = data['activities'] if isinstance(data['activities'], list) else [data['activities']]
            tf = content.text_frame
            for activity in activities:
                if isinstance(activity, dict):
                    name = activity.get('name', 'Activity')
                    desc = activity.get('description', '')
                    p = tf.add_paragraph()
                    p.text = f"{name}: {desc}"
                    p.level = 0
                else:
                    p = tf.add_paragraph()
                    p.text = str(activity)
                    p.level = 0
    
    def _format_generic_data_for_powerpoint(self, prs: Presentation, data: Dict[str, Any], layout):
        """Format generic data for PowerPoint."""
        slide = prs.slides.add_slide(layout)
        title = slide.shapes.title
        content = slide.placeholders[1]
        title.text = "Widget Data"
        
        tf = content.text_frame
        if isinstance(data, dict):
            for key, value in data.items():
                if value and key not in ['id', 'widget_id']:
                    p = tf.add_paragraph()
                    p.text = f"{key.replace('_', ' ').title()}: {str(value)[:100]}"
                    p.level = 0
        elif isinstance(data, list):
            for item in data[:10]:  # Limit to 10 items
                p = tf.add_paragraph()
                p.text = str(item)[:100]
                p.level = 0
    
    def create_presentation_from_slides(
        self,
        presentation_title: str,
        slides: List[Dict[str, Any]],
        subtitle: Optional[str] = None
    ) -> bytes:
        """
        Create a PowerPoint presentation from a list of slides.
        
        Args:
            presentation_title: Title of the presentation
            slides: List of slide dictionaries, each with:
                - title: Slide title (required)
                - content: Slide content (can be string or list of bullet points)
                - notes: Optional speaker notes
            subtitle: Optional subtitle for title slide
            
        Returns:
            PowerPoint file as bytes
        """
        try:
            prs = Presentation()
            prs.slide_width = PptInches(10)
            prs.slide_height = PptInches(7.5)
            
            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle_placeholder = slide.placeholders[1]
            
            title.text = presentation_title
            if subtitle:
                subtitle_placeholder.text = subtitle
            else:
                subtitle_placeholder.text = f"Created: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S UTC')}"
            
            # Content slide layout
            bullet_slide_layout = prs.slide_layouts[1]
            
            # Add content slides
            for slide_data in slides:
                slide = prs.slides.add_slide(bullet_slide_layout)
                title_shape = slide.shapes.title
                content_placeholder = slide.placeholders[1]
                
                # Set title
                title_shape.text = slide_data.get('title', 'Untitled Slide')
                
                # Add image if provided
                if slide_data.get('image_url') or slide_data.get('image_base64'):
                    try:
                        image_bytes = None
                        if slide_data.get('image_base64'):
                            image_bytes = base64.b64decode(slide_data.get('image_base64'))
                        elif slide_data.get('image_url'):
                            response = requests.get(slide_data.get('image_url'))
                            image_bytes = response.content
                        
                        if image_bytes:
                            # Add image to slide (right side)
                            left = PptInches(5.5)
                            top = PptInches(1.5)
                            width = PptInches(4)
                            height = PptInches(3)
                            slide.shapes.add_picture(io.BytesIO(image_bytes), left, top, width, height)
                    except Exception as e:
                        self.logger.warning(f"Could not add image to slide: {str(e)}")
                
                # Set content
                tf = content_placeholder.text_frame
                tf.text = ""  # Clear default text
                
                content = slide_data.get('content', '')
                if isinstance(content, list):
                    # Multiple bullet points
                    for i, bullet in enumerate(content):
                        if i == 0:
                            p = tf.paragraphs[0]
                            p.text = str(bullet)
                            p.level = 0
                        else:
                            p = tf.add_paragraph()
                            p.text = str(bullet)
                            p.level = 0
                elif isinstance(content, str):
                    # Single paragraph or multi-line text
                    lines = content.split('\n')
                    for i, line in enumerate(lines):
                        if line.strip():
                            if i == 0:
                                p = tf.paragraphs[0]
                                p.text = line.strip()
                                p.level = 0
                            else:
                                p = tf.add_paragraph()
                                p.text = line.strip()
                                p.level = 0
                
                # Add video/web links if provided
                if slide_data.get('video_links'):
                    for video_link in slide_data.get('video_links', []):
                        p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = f"Video: {video_link.get('title', 'Watch Video')}"
                        run.hyperlink.address = video_link.get('url', '')
                
                if slide_data.get('web_links'):
                    for web_link in slide_data.get('web_links', []):
                        p = tf.add_paragraph()
                        run = p.add_run()
                        run.text = f"Resource: {web_link.get('title', 'View Resource')}"
                        run.hyperlink.address = web_link.get('url', '')
                
                # Add speaker notes if provided
                if slide_data.get('notes'):
                    notes_slide = slide.notes_slide
                    notes_text_frame = notes_slide.notes_text_frame
                    notes_text_frame.text = slide_data.get('notes')
            
            # Save to bytes
            buffer = io.BytesIO()
            prs.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            self.logger.error(f"Error creating PowerPoint presentation: {str(e)}")
            raise
    
    def create_word_document_from_content(
        self,
        document_title: str,
        content: List[Dict[str, Any]],
        subtitle: Optional[str] = None
    ) -> bytes:
        """
        Create a Word document from content sections.
        
        Args:
            document_title: Title of the document
            content: List of content sections, each with:
                - heading: Section heading (required)
                - paragraphs: List of paragraph text (required)
                - bullet_points: Optional list of bullet points
            subtitle: Optional subtitle for the document
            
        Returns:
            Word document file as bytes
        """
        try:
            doc = Document()
            
            # Title
            title = doc.add_heading(document_title, 0)
            title.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            # Subtitle or date
            if subtitle:
                doc.add_paragraph(subtitle).alignment = WD_ALIGN_PARAGRAPH.CENTER
            else:
                date_para = doc.add_paragraph(f"Created: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S UTC')}")
                date_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
            
            doc.add_paragraph()  # Blank line
            
            # Add content sections
            for section in content:
                # Section heading
                heading = doc.add_heading(section.get('heading', 'Section'), level=1)
                
                # Add image if provided
                has_image_url = bool(section.get('image_url'))
                has_image_base64 = bool(section.get('image_base64'))
                if has_image_url or has_image_base64:
                    self.logger.info(f"📷 Found image in section '{section.get('heading', 'Section')}': image_url={has_image_url}, image_base64={has_image_base64} (length: {len(section.get('image_base64', '')) if section.get('image_base64') else 0})")
                    try:
                        image_bytes = None
                        if section.get('image_base64'):
                            # Handle base64 string - remove data URL prefix if present
                            base64_str = section.get('image_base64')
                            if isinstance(base64_str, str):
                                # Remove data URL prefix if present (e.g., "data:image/png;base64,")
                                if ',' in base64_str:
                                    base64_str = base64_str.split(',', 1)[1]
                                # Remove any whitespace
                                base64_str = base64_str.strip()
                                # Add padding if needed
                                missing_padding = len(base64_str) % 4
                                if missing_padding:
                                    base64_str += '=' * (4 - missing_padding)
                                try:
                                    image_bytes = base64.b64decode(base64_str)
                                    self.logger.info(f"✅ Successfully decoded base64 image, size: {len(image_bytes)} bytes")
                                except Exception as decode_error:
                                    self.logger.warning(f"Base64 decode error: {str(decode_error)}, attempting without padding fix")
                                    # Try without padding fix as fallback
                                    try:
                                        image_bytes = base64.b64decode(base64_str, validate=False)
                                        self.logger.info(f"✅ Successfully decoded base64 image (without validation), size: {len(image_bytes)} bytes")
                                    except Exception as decode_error2:
                                        self.logger.error(f"❌ Failed to decode base64 image even without validation: {str(decode_error2)}")
                                        image_bytes = None
                        elif section.get('image_url'):
                            try:
                                response = requests.get(section.get('image_url'), timeout=10)
                                response.raise_for_status()
                                image_bytes = response.content
                                self.logger.info(f"✅ Successfully downloaded image from URL, size: {len(image_bytes)} bytes")
                            except Exception as url_error:
                                self.logger.error(f"❌ Failed to download image from URL: {str(url_error)}")
                                image_bytes = None
                        
                        if image_bytes:
                            # Add image to document
                            image_stream = io.BytesIO(image_bytes)
                            doc.add_picture(image_stream, width=Inches(5))
                            doc.add_paragraph()  # Blank line after image
                            self.logger.info(f"✅ Successfully added image to Word document")
                        else:
                            self.logger.warning(f"⚠️ No image bytes available to add to document (image_base64: {bool(section.get('image_base64'))}, image_url: {bool(section.get('image_url'))})")
                    except Exception as e:
                        error_msg = str(e) if str(e) else repr(e)
                        self.logger.warning(f"Could not add image to document: {error_msg}", exc_info=True)
                        # Continue without the image rather than failing
                
                # Paragraphs
                if section.get('paragraphs'):
                    for para_text in section.get('paragraphs', []):
                        if para_text:
                            para = doc.add_paragraph(para_text)
                            # Check if paragraph contains a URL and make it a hyperlink
                            import re
                            urls = re.findall(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', para_text)
                            if urls:
                                para.clear()
                                parts = re.split(r'(http[s]?://[^\s]+)', para_text)
                                for part in parts:
                                    if part.startswith('http'):
                                        para.add_run(part).hyperlink.address = part
                                    else:
                                        para.add_run(part)
                
                # Bullet points
                if section.get('bullet_points'):
                    for bullet in section.get('bullet_points', []):
                        if bullet:
                            para = doc.add_paragraph(bullet, style='List Bullet')
                
                # Add video/web links if provided
                if section.get('video_links'):
                    doc.add_paragraph('Video Resources:', style='Heading 3')
                    for video_link in section.get('video_links', []):
                        para = doc.add_paragraph()
                        hyperlink = para.add_run(video_link.get('title', 'Watch Video'))
                        hyperlink.hyperlink.address = video_link.get('url', '')
                
                if section.get('web_links'):
                    doc.add_paragraph('Web Resources:', style='Heading 3')
                    for web_link in section.get('web_links', []):
                        para = doc.add_paragraph()
                        hyperlink = para.add_run(web_link.get('title', 'View Resource'))
                        hyperlink.hyperlink.address = web_link.get('url', '')
                
                doc.add_paragraph()  # Blank line between sections
            
            # Save to bytes
            buffer = io.BytesIO()
            doc.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            self.logger.error(f"Error creating Word document: {str(e)}")
            raise
    
    def create_pdf_from_content(
        self,
        document_title: str,
        content: List[Dict[str, Any]],
        subtitle: Optional[str] = None
    ) -> bytes:
        """
        Create a PDF document from content sections.
        
        Args:
            document_title: Title of the document
            content: List of content sections, each with:
                - heading: Section heading (required)
                - paragraphs: List of paragraph text (required)
                - bullet_points: Optional list of bullet points
            subtitle: Optional subtitle for the document
            
        Returns:
            PDF file as bytes
        """
        try:
            buffer = io.BytesIO()
            doc = SimpleDocTemplate(buffer, pagesize=letter)
            story = []
            styles = getSampleStyleSheet()
            
            # Title style
            title_style = ParagraphStyle(
                'DocumentTitle',
                parent=styles['Heading1'],
                fontSize=20,
                textColor=colors.HexColor('#1a1a1a'),
                spaceAfter=20,
                alignment=TA_CENTER
            )
            
            # Add title
            story.append(Paragraph(document_title, title_style))
            
            # Subtitle or date
            if subtitle:
                story.append(Paragraph(subtitle, styles['Normal']))
            else:
                story.append(Paragraph(f"Created: {datetime.utcnow().strftime('%Y-%m-%d %H:%M:%S UTC')}", styles['Normal']))
            
            story.append(Spacer(1, 0.3*inch))
            
            # Add content sections
            for section in content:
                # Section heading
                heading = section.get('heading', 'Section')
                story.append(Paragraph(heading, styles['Heading2']))
                story.append(Spacer(1, 0.1*inch))
                
                # Add image if provided
                if section.get('image_url') or section.get('image_base64'):
                    try:
                        image_bytes = None
                        if section.get('image_base64'):
                            image_bytes = base64.b64decode(section.get('image_base64'))
                        elif section.get('image_url'):
                            response = requests.get(section.get('image_url'))
                            image_bytes = response.content
                        
                        if image_bytes:
                            img = ReportLabImage(io.BytesIO(image_bytes), width=5*inch, height=3*inch)
                            story.append(img)
                            story.append(Spacer(1, 0.1*inch))
                    except Exception as e:
                        self.logger.warning(f"Could not add image to PDF: {str(e)}")
                
                # Paragraphs
                if section.get('paragraphs'):
                    for para_text in section.get('paragraphs', []):
                        if para_text:
                            # Check for URLs and format as hyperlinks
                            import re
                            urls = re.findall(r'http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+', para_text)
                            if urls:
                                # Split text and create hyperlinks
                                parts = re.split(r'(http[s]?://[^\s]+)', para_text)
                                para_parts = []
                                for part in parts:
                                    if part.startswith('http'):
                                        from reportlab.lib.styles import getSampleStyleSheet
                                        link_style = ParagraphStyle('Link', parent=styles['Normal'], textColor=colors.HexColor('#0066cc'), underline=1)
                                        para_parts.append(Paragraph(f'<a href="{part}" color="blue">{part}</a>', link_style))
                                    else:
                                        para_parts.append(Paragraph(part, styles['Normal']))
                                story.extend(para_parts)
                            else:
                                story.append(Paragraph(para_text, styles['Normal']))
                            story.append(Spacer(1, 0.05*inch))
                
                # Bullet points
                if section.get('bullet_points'):
                    for bullet in section.get('bullet_points', []):
                        if bullet:
                            story.append(Paragraph(f"• {bullet}", styles['Normal']))
                            story.append(Spacer(1, 0.03*inch))
                
                # Add video/web links if provided
                if section.get('video_links'):
                    story.append(Paragraph('Video Resources:', styles['Heading3']))
                    for video_link in section.get('video_links', []):
                        link_style = ParagraphStyle('Link', parent=styles['Normal'], textColor=colors.HexColor('#0066cc'), underline=1)
                        story.append(Paragraph(f'<a href="{video_link.get("url", "")}" color="blue">{video_link.get("title", "Watch Video")}</a>', link_style))
                        story.append(Spacer(1, 0.05*inch))
                
                if section.get('web_links'):
                    story.append(Paragraph('Web Resources:', styles['Heading3']))
                    for web_link in section.get('web_links', []):
                        link_style = ParagraphStyle('Link', parent=styles['Normal'], textColor=colors.HexColor('#0066cc'), underline=1)
                        story.append(Paragraph(f'<a href="{web_link.get("url", "")}" color="blue">{web_link.get("title", "View Resource")}</a>', link_style))
                        story.append(Spacer(1, 0.05*inch))
                
                story.append(Spacer(1, 0.2*inch))
            
            # Build PDF
            doc.build(story)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            self.logger.error(f"Error creating PDF document: {str(e)}")
            raise
    
    def create_excel_spreadsheet_from_data(
        self,
        spreadsheet_title: str,
        sheets: List[Dict[str, Any]],
        subtitle: Optional[str] = None
    ) -> bytes:
        """
        Create an Excel spreadsheet from data.
        
        Args:
            spreadsheet_title: Title of the spreadsheet
            sheets: List of sheet data, each with:
                - name: Sheet name (required)
                - headers: List of column headers (required)
                - rows: List of row data (list of lists or list of dicts) (required)
                - summary: Optional summary text for the sheet
            subtitle: Optional subtitle
            
        Returns:
            Excel file as bytes
        """
        try:
            wb = Workbook()
            
            # Remove default sheet if we have custom sheets
            if len(sheets) > 0:
                wb.remove(wb.active)
            
            # Header style
            header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
            header_font = Font(bold=True, color="FFFFFF", size=12)
            border = Border(
                left=Side(style='thin'),
                right=Side(style='thin'),
                top=Side(style='thin'),
                bottom=Side(style='thin')
            )
            
            # Create sheets
            for sheet_data in sheets:
                sheet_name = sheet_data.get('name', 'Sheet1')[:31]  # Excel limit
                ws = wb.create_sheet(title=sheet_name)
                
                # Title row
                headers = sheet_data.get('headers', [])
                num_cols = len(headers)
                if num_cols > 0:
                    ws.merge_cells(f'A1:{get_column_letter(num_cols)}1')
                    title_cell = ws['A1']
                    title_cell.value = spreadsheet_title
                    title_cell.font = Font(bold=True, size=14)
                    title_cell.alignment = Alignment(horizontal='center', vertical='center')
                    
                    if subtitle:
                        ws.merge_cells(f'A2:{get_column_letter(num_cols)}2')
                        subtitle_cell = ws['A2']
                        subtitle_cell.value = subtitle
                        subtitle_cell.alignment = Alignment(horizontal='center', vertical='center')
                    
                    # Headers
                    header_row = 3 if subtitle else 2
                    for col_idx, header in enumerate(headers, 1):
                        cell = ws.cell(row=header_row, column=col_idx)
                        cell.value = header
                        cell.fill = header_fill
                        cell.font = header_font
                        cell.border = border
                        cell.alignment = Alignment(horizontal='center', vertical='center')
                    
                    # Data rows
                    rows = sheet_data.get('rows', [])
                    for row_idx, row_data in enumerate(rows, start=header_row + 1):
                        if isinstance(row_data, dict):
                            # Map dict keys to columns
                            for col_idx, header in enumerate(headers, 1):
                                cell = ws.cell(row=row_idx, column=col_idx)
                                cell.value = row_data.get(header, '')
                                cell.border = border
                        elif isinstance(row_data, list):
                            # List of values
                            for col_idx, value in enumerate(row_data[:num_cols], 1):
                                cell = ws.cell(row=row_idx, column=col_idx)
                                cell.value = value
                                cell.border = border
                    
                    # Auto-adjust column widths
                    for column in ws.columns:
                        max_length = 0
                        column_letter = get_column_letter(column[0].column)
                        for cell in column:
                            try:
                                if cell.value and len(str(cell.value)) > max_length:
                                    max_length = len(str(cell.value))
                            except:
                                pass
                        adjusted_width = min(max_length + 2, 50)
                        ws.column_dimensions[column_letter].width = adjusted_width
                
                # Summary if provided
                if sheet_data.get('summary'):
                    summary_row = header_row + len(rows) + 2
                    ws.merge_cells(f'A{summary_row}:{get_column_letter(num_cols)}{summary_row}')
                    summary_cell = ws[f'A{summary_row}']
                    summary_cell.value = sheet_data.get('summary')
                    summary_cell.font = Font(italic=True)
            
            # If no sheets provided, create a default one
            if len(sheets) == 0:
                ws = wb.active
                ws.title = spreadsheet_title[:31]
                ws['A1'] = spreadsheet_title
            
            # Save to bytes
            buffer = io.BytesIO()
            wb.save(buffer)
            buffer.seek(0)
            return buffer.getvalue()
            
        except Exception as e:
            self.logger.error(f"Error creating Excel spreadsheet: {str(e)}")
            raise

